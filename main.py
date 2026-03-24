# ╔══════════════════════════════════════════════════════════════════╗
# ║         ADVANCED RAG CHATBOT — FastAPI + Ollama + ChromaDB       ║
# ║         Multi-doc · Streaming · Memory · Source Citations        ║
# ╚══════════════════════════════════════════════════════════════════╝

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import HTMLResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import ollama
import chromadb
from chromadb.utils import embedding_functions
import uuid
import os
import io
import json
import re
import zipfile
import tempfile
import hashlib
from pathlib import Path
from typing import Optional
from datetime import datetime

# ── Optional parsers (install what you need) ──────────────────────
try:
    from pypdf import PdfReader
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from langchain.text_splitter import RecursiveCharacterTextSplitter


# ── CONFIG ────────────────────────────────────────────────────────
EMBED_MODEL   = "nomic-embed-text"
DEFAULT_MODEL = "llama3"
DB_PATH       = "./chroma_db"
UPLOAD_DIR    = "./uploads"
CHUNK_SIZE    = 600
CHUNK_OVERLAP = 80
TOP_K         = 5

os.makedirs(UPLOAD_DIR, exist_ok=True)

app = FastAPI(title="Advanced RAG Chatbot", version="2.0")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── CHROMADB ──────────────────────────────────────────────────────
chroma_client = chromadb.PersistentClient(path=DB_PATH)

# In-memory session store: session_id → {collection, history, docs}
sessions: dict = {}


# ══════════════════════════════════════════════════════════════════
#  SESSION MANAGEMENT
# ══════════════════════════════════════════════════════════════════

def get_or_create_session(session_id: str) -> dict:
    if session_id not in sessions:
        collection = chroma_client.get_or_create_collection(
            name=f"session_{session_id}",
            metadata={"hnsw:space": "cosine"}
        )
        sessions[session_id] = {
            "collection": collection,
            "history": [],       # [{role, content}]
            "documents": {},     # {doc_id: {name, chunks, uploaded_at}}
        }
    return sessions[session_id]


# ══════════════════════════════════════════════════════════════════
#  TEXT EXTRACTION
# ══════════════════════════════════════════════════════════════════

def extract_text_pdf(data: bytes, filename: str) -> list[dict]:
    """Returns list of {text, page, source}"""
    if not HAS_PDF:
        raise HTTPException(400, "pypdf not installed. Run: pip install pypdf")
    pages = []
    reader = PdfReader(io.BytesIO(data))
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({"text": text, "page": i + 1, "source": filename})
    return pages


def extract_text_docx(data: bytes, filename: str) -> list[dict]:
    if not HAS_DOCX:
        raise HTTPException(400, "python-docx not installed. Run: pip install python-docx")
    doc = DocxDocument(io.BytesIO(data))
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"text": full_text, "page": 1, "source": filename}]


def extract_text_txt(data: bytes, filename: str) -> list[dict]:
    text = data.decode("utf-8", errors="ignore")
    return [{"text": text, "page": 1, "source": filename}]


def extract_text_csv(data: bytes, filename: str) -> list[dict]:
    if not HAS_PANDAS:
        text = data.decode("utf-8", errors="ignore")
        return [{"text": text, "page": 1, "source": filename}]
    df = pd.read_csv(io.BytesIO(data))
    # Convert each row to readable text
    rows = []
    for i, row in df.iterrows():
        row_text = " | ".join(f"{col}: {val}" for col, val in row.items())
        rows.append(row_text)
    full_text = "\n".join(rows)
    return [{"text": full_text, "page": 1, "source": filename}]


def extract_text_pptx(data: bytes, filename: str) -> list[dict]:
    if not HAS_PPTX:
        raise HTTPException(400, "python-pptx not installed. Run: pip install python-pptx")
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        if texts:
            slides.append({"text": "\n".join(texts), "page": i + 1, "source": filename})
    return slides


def extract_text(data: bytes, filename: str) -> list[dict]:
    ext = Path(filename).suffix.lower()
    extractors = {
        ".pdf":  extract_text_pdf,
        ".docx": extract_text_docx,
        ".txt":  extract_text_txt,
        ".csv":  extract_text_csv,
        ".pptx": extract_text_pptx,
    }
    if ext not in extractors:
        raise HTTPException(400, f"Unsupported file type: {ext}")
    return extractors[ext](data, filename)


# ══════════════════════════════════════════════════════════════════
#  CHUNKING & EMBEDDING
# ══════════════════════════════════════════════════════════════════

splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    separators=["\n\n", "\n", ". ", " ", ""]
)


def chunk_pages(pages: list[dict]) -> list[dict]:
    """Split pages into smaller chunks, keeping metadata."""
    chunks = []
    for page in pages:
        splits = splitter.split_text(page["text"])
        for j, split in enumerate(splits):
            if split.strip():
                chunks.append({
                    "text": split.strip(),
                    "page": page["page"],
                    "source": page["source"],
                    "chunk_index": j,
                })
    return chunks


def get_embedding(text: str) -> list[float]:
    resp = ollama.embeddings(model=EMBED_MODEL, prompt=text)
    return resp["embedding"]


def dedup_chunks(chunks: list[dict]) -> list[dict]:
    """Remove duplicate chunks by content hash."""
    seen = set()
    unique = []
    for c in chunks:
        h = hashlib.md5(c["text"].encode()).hexdigest()
        if h not in seen:
            seen.add(h)
            unique.append(c)
    return unique


# ══════════════════════════════════════════════════════════════════
#  INGESTION
# ══════════════════════════════════════════════════════════════════

def ingest_document(session_id: str, filename: str, data: bytes) -> dict:
    session = get_or_create_session(session_id)
    collection = session["collection"]

    # Extract text
    pages = extract_text(data, filename)
    if not pages:
        raise HTTPException(400, f"Could not extract text from {filename}")

    # Chunk
    chunks = chunk_pages(pages)
    chunks = dedup_chunks(chunks)

    # Embed and store
    doc_id = str(uuid.uuid4())[:8]
    ids, embeddings, documents, metadatas = [], [], [], []

    for i, chunk in enumerate(chunks):
        chunk_id = f"{doc_id}_{i}"
        emb = get_embedding(chunk["text"])
        ids.append(chunk_id)
        embeddings.append(emb)
        documents.append(chunk["text"])
        metadatas.append({
            "source": chunk["source"],
            "page": chunk["page"],
            "doc_id": doc_id,
            "chunk_index": i,
        })

    # Batch add
    batch = 100
    for i in range(0, len(ids), batch):
        collection.add(
            ids=ids[i:i+batch],
            embeddings=embeddings[i:i+batch],
            documents=documents[i:i+batch],
            metadatas=metadatas[i:i+batch],
        )

    # Record in session
    session["documents"][doc_id] = {
        "name": filename,
        "chunks": len(chunks),
        "pages": len(pages),
        "uploaded_at": datetime.now().isoformat(),
        "doc_id": doc_id,
    }

    return {"doc_id": doc_id, "chunks": len(chunks), "pages": len(pages)}


# ══════════════════════════════════════════════════════════════════
#  HYBRID RETRIEVAL
# ══════════════════════════════════════════════════════════════════

def keyword_score(text: str, query: str) -> float:
    """Simple BM25-style keyword overlap score."""
    query_words = set(re.findall(r'\w+', query.lower()))
    text_words  = set(re.findall(r'\w+', text.lower()))
    if not query_words:
        return 0.0
    overlap = len(query_words & text_words)
    return overlap / len(query_words)


def hybrid_retrieve(session_id: str, query: str, k: int = TOP_K) -> list[dict]:
    session = get_or_create_session(session_id)
    collection = session["collection"]

    if collection.count() == 0:
        return []

    # Semantic search — get more candidates then re-rank
    q_emb = get_embedding(query)
    n = min(k * 3, collection.count())
    results = collection.query(
        query_embeddings=[q_emb],
        n_results=n,
        include=["documents", "metadatas", "distances"]
    )

    docs      = results["documents"][0]
    metas     = results["metadatas"][0]
    distances = results["distances"][0]

    # Hybrid score = semantic (1 - cosine distance) + keyword
    scored = []
    seen_texts = set()
    for doc, meta, dist in zip(docs, metas, distances):
        if doc in seen_texts:
            continue
        seen_texts.add(doc)
        semantic  = 1.0 - dist
        keyword   = keyword_score(doc, query)
        combined  = 0.7 * semantic + 0.3 * keyword
        scored.append({
            "text":    doc,
            "source":  meta.get("source", "unknown"),
            "page":    meta.get("page", "?"),
            "doc_id":  meta.get("doc_id", ""),
            "score":   combined,
        })

    scored.sort(key=lambda x: x["score"], reverse=True)
    return scored[:k]


# ══════════════════════════════════════════════════════════════════
#  STREAMING RAG RESPONSE
# ══════════════════════════════════════════════════════════════════

def build_prompt(query: str, chunks: list[dict], history: list[dict]) -> list[dict]:
    # Build context block
    context_parts = []
    for i, c in enumerate(chunks):
        context_parts.append(
            f"[Source {i+1}: {c['source']}, page {c['page']}]\n{c['text']}"
        )
    context = "\n\n---\n\n".join(context_parts)

    system_msg = (
        "You are a helpful document assistant. "
        "Answer questions ONLY using the provided context below. "
        "Always cite your sources by mentioning the document name and page number. "
        "If the answer is not found in the context, say: "
        "'I could not find information about this in the uploaded documents.' "
        "Do not make up information.\n\n"
        f"CONTEXT:\n{context}"
    )

    messages = [{"role": "system", "content": system_msg}]

    # Add last 6 history turns (3 exchanges) for memory
    for turn in history[-6:]:
        messages.append(turn)

    messages.append({"role": "user", "content": query})
    return messages


def stream_response(session_id: str, query: str, model: str):
    session  = get_or_create_session(session_id)
    history  = session["history"]
    chunks   = hybrid_retrieve(session_id, query)

    # Build source citations header
    if chunks:
        sources = []
        seen = set()
        for c in chunks:
            key = f"{c['source']}::{c['page']}"
            if key not in seen:
                seen.add(key)
                sources.append({"source": c["source"], "page": c["page"]})

        # Stream sources as first SSE event
        yield f"data: {json.dumps({'type': 'sources', 'sources': sources})}\n\n"
    else:
        yield f"data: {json.dumps({'type': 'sources', 'sources': []})}\n\n"

    messages = build_prompt(query, chunks, history)

    # Stream tokens
    full_response = ""
    try:
        stream = ollama.chat(model=model, messages=messages, stream=True)
        for part in stream:
            token = part["message"]["content"]
            full_response += token
            yield f"data: {json.dumps({'type': 'token', 'token': token})}\n\n"
    except Exception as e:
        err = f"Error: {str(e)}"
        yield f"data: {json.dumps({'type': 'token', 'token': err})}\n\n"
        full_response = err

    # Save to history
    session["history"].append({"role": "user",      "content": query})
    session["history"].append({"role": "assistant",  "content": full_response})

    yield f"data: {json.dumps({'type': 'done'})}\n\n"


# ══════════════════════════════════════════════════════════════════
#  REQUEST MODELS
# ══════════════════════════════════════════════════════════════════

class ChatRequest(BaseModel):
    session_id: str
    question: str
    model: str = DEFAULT_MODEL


class SessionRequest(BaseModel):
    session_id: Optional[str] = None


# ══════════════════════════════════════════════════════════════════
#  API ROUTES
# ══════════════════════════════════════════════════════════════════

@app.post("/session/new")
def new_session():
    sid = str(uuid.uuid4())[:12]
    get_or_create_session(sid)
    return {"session_id": sid}


@app.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    session_id: str  = Form(...),
):
    data = await file.read()
    if len(data) > 50 * 1024 * 1024:  # 50MB limit
        raise HTTPException(400, "File too large (max 50MB)")

    result = ingest_document(session_id, file.filename, data)
    return {
        "success": True,
        "filename": file.filename,
        "doc_id":   result["doc_id"],
        "chunks":   result["chunks"],
        "pages":    result["pages"],
    }


@app.post("/upload/zip")
async def upload_zip(
    file: UploadFile = File(...),
    session_id: str  = Form(...),
):
    data = await file.read()
    results = []
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        for name in z.namelist():
            ext = Path(name).suffix.lower()
            if ext in {".pdf", ".docx", ".txt", ".csv", ".pptx"}:
                file_data = z.read(name)
                try:
                    r = ingest_document(session_id, Path(name).name, file_data)
                    results.append({"file": name, **r})
                except Exception as e:
                    results.append({"file": name, "error": str(e)})
    return {"success": True, "files": results}


@app.get("/chat/stream")
async def chat_stream(session_id: str, question: str, model: str = DEFAULT_MODEL):
    return StreamingResponse(
        stream_response(session_id, question, model),
        media_type="text/event-stream",
        headers={
            "Cache-Control":               "no-cache",
            "X-Accel-Buffering":           "no",
            "Access-Control-Allow-Origin": "*",
        }
    )


@app.get("/session/{session_id}/documents")
def get_documents(session_id: str):
    session = get_or_create_session(session_id)
    return {"documents": list(session["documents"].values())}


@app.get("/session/{session_id}/history")
def get_history(session_id: str):
    session = get_or_create_session(session_id)
    return {"history": session["history"]}


@app.delete("/session/{session_id}/history")
def clear_history(session_id: str):
    session = get_or_create_session(session_id)
    session["history"] = []
    return {"success": True}


@app.delete("/session/{session_id}/document/{doc_id}")
def delete_document(session_id: str, doc_id: str):
    session = get_or_create_session(session_id)
    collection = session["collection"]
    # Delete all chunks belonging to this doc
    results = collection.get(where={"doc_id": doc_id})
    if results["ids"]:
        collection.delete(ids=results["ids"])
    session["documents"].pop(doc_id, None)
    return {"success": True}


@app.get("/health")
def health():
    return {
        "status":       "ok",
        "embed_model":  EMBED_MODEL,
        "sessions":     len(sessions),
    }


@app.get("/models")
def list_models():
    try:
        resp = ollama.list()
        names = [m["name"] for m in resp.get("models", [])]
        return {"models": names}
    except Exception:
        return {"models": ["llama3", "mistral", "phi3", "gemma"]}


# ══════════════════════════════════════════════════════════════════
#  SERVE FRONTEND
# ══════════════════════════════════════════════════════════════════

@app.get("/", response_class=HTMLResponse)
def serve_ui():
    ui_path = Path(__file__).parent / "chat_ui.html"
    if ui_path.exists():
        return ui_path.read_text(encoding="utf-8")

    return "<h2>Place chat_ui.html in the same folder as main.py</h2>"


# ══════════════════════════════════════════════════════════════════
#  ENTRY POINT
# ══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
